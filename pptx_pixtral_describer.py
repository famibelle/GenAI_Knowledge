#!/usr/bin/env python3
"""
Génère une description détaillée de chaque slide d'un fichier PPTX en utilisant le modèle vision Pixtral (Mistral multimodal)
Résultat: un fichier Excel avec colonnes:
  - Slide Number
  - Slide Title
  - description détaillée de la slides

Prérequis:
  pip install python-pptx pandas transformers accelerate bitsandbytes pillow
  (Optionnel pour export images via LibreOffice): sudo apt-get install libreoffice-common

Si LibreOffice est installé, les slides sont rendues en PNG haute fidélité.
Sinon, on tombera en mode fallback: on reconstruit une 'image' blanche contenant le texte concaténé (qualité moindre).

Usage:
  python pptx_pixtral_describer.py --pptx cours.pptx --out descriptions.xlsx

Options principales:
  --model mistral-community/pixtral-12b
  --quant 4bit|8bit|fp16|auto   (par défaut: 4bit pour GPU T4)

Limitation: le rendu fidèle nécessite une conversion PPTX->PNG (LibreOffice ou autre outil). Le script tente automatiquement.
"""
from __future__ import annotations
import argparse
import os
import sys
import subprocess
import glob
import tempfile
from dataclasses import dataclass
from typing import List, Optional

import pandas as pd
from pptx import Presentation
from PIL import Image, ImageDraw, ImageFont
import torch
import math

# ----------------------------------------------------------
# Chargement modèle vision (Pixtral) avec quantization
# ----------------------------------------------------------
from transformers import AutoProcessor, AutoModelForCausalLM
try:
    from transformers import BitsAndBytesConfig
except ImportError:  # ancienne version transformers
    BitsAndBytesConfig = None


def _print_env():
    try:
        import bitsandbytes as bnb  # noqa
        bnb_version = bnb.__version__
    except Exception:
        bnb_version = 'non installé'
    gpu = torch.cuda.is_available()
    gpu_name = torch.cuda.get_device_name(0) if gpu else 'CPU'
    cap = torch.cuda.get_device_capability(0) if gpu else ('-','-')
    print(f"🔍 Torch {torch.__version__} | bitsandbytes {bnb_version} | GPU: {gpu_name} cc{cap}")


def load_pixtral(model_name: str, quant: str, attn_impl: str = 'eager', trust_remote_code: bool = True,
                 force_cpu: bool = False, offload_dir: Optional[str] = None, max_gpu_mem: Optional[str] = None,
                 fast_processor: bool = False):
    quant = (quant or 'auto').lower()
    gpu = torch.cuda.is_available() and not force_cpu
    _print_env()

    load_kwargs = {
        "device_map": "cpu" if force_cpu else "auto",
        "trust_remote_code": trust_remote_code,
    }
    if fast_processor:
        os.environ.setdefault('HF_USE_FAST_PROCESSOR', '1')
    # Allocation config pour limiter fragmentation
    os.environ.setdefault('PYTORCH_CUDA_ALLOC_CONF', 'expandable_segments:True')
    if offload_dir:
        os.makedirs(offload_dir, exist_ok=True)
        load_kwargs['offload_folder'] = offload_dir
    if max_gpu_mem and gpu:
        load_kwargs['max_memory'] = {0: max_gpu_mem, 'cpu': '48GiB'}
    # Attention impl (certains flash_attn causent des segfaults selon versions)
    if attn_impl:
        os.environ.setdefault('ATTN_IMPLEMENTATION', attn_impl)

    def build_bnb_config(mode: str):
        if BitsAndBytesConfig is None:
            return None
        if mode == '4bit':
            return BitsAndBytesConfig(
                load_in_4bit=True,
                bnb_4bit_quant_type='nf4',
                bnb_4bit_compute_dtype=torch.bfloat16 if gpu and torch.cuda.is_bf16_supported() else torch.float16,
                bnb_4bit_use_double_quant=True,
            )
        if mode == '8bit':
            return BitsAndBytesConfig(load_in_8bit=True)
        return None

    attempt_order = []
    if quant in ('4bit','4'):
        attempt_order = ['4bit']
    elif quant in ('8bit','8'):
        attempt_order = ['8bit']
    elif quant in ('fp16','half','16'):
        attempt_order = ['fp16']
    else:  # auto
        attempt_order = ['4bit','8bit','fp16']

    last_error = None
    processor = AutoProcessor.from_pretrained(model_name, trust_remote_code=trust_remote_code, use_fast=fast_processor)
    for mode in attempt_order:
        local_kwargs = dict(load_kwargs)
        quant_cfg = None
        if mode == '4bit':
            quant_cfg = build_bnb_config('4bit') if gpu else None
        elif mode == '8bit':
            quant_cfg = build_bnb_config('8bit') if gpu else None
        elif mode == 'fp16' and gpu:
            local_kwargs['torch_dtype'] = torch.float16
        if quant_cfg:
            local_kwargs['quantization_config'] = quant_cfg
        print(f"⚙️  Tentative chargement mode={mode} ...")
        try:
            model = AutoModelForCausalLM.from_pretrained(model_name, **local_kwargs)
            print(f"✅ Modèle chargé ({mode})")
            return processor, model
        except Exception as e:
            print(f"⚠️  Échec mode {mode}: {e}")
            last_error = e

    # Ultimate fallback: CPU bfloat16/float32 offload
    print("↪️  Fallback CPU + offload progressif ...")
    cpu_kwargs = {"device_map": "auto", "trust_remote_code": trust_remote_code}
    if torch.cuda.is_available() and not force_cpu:
        # mémoire GPU disponible pour partition
        total = torch.cuda.get_device_properties(0).total_memory
        reserve = int(1.2 * 1024**3)
        allowance = max(total - reserve, int(0.5 * 1024**3))
        allowance_gib = f"{math.floor(allowance/1024**3)}GiB"
        cpu_kwargs['max_memory'] = {0: allowance_gib, 'cpu': '96GiB'}
    if offload_dir:
        cpu_kwargs['offload_folder'] = offload_dir
    try:
        model = AutoModelForCausalLM.from_pretrained(model_name, **cpu_kwargs)
        print("✅ Modèle chargé (fallback CPU/offload)")
        return processor, model
    except Exception as e2:
        print(f"❌ Impossible de charger le modèle après fallbacks: {e2}\nDernière erreur quant: {last_error}")
        raise

# ----------------------------------------------------------
# Export slides en images (LibreOffice)
# ----------------------------------------------------------

def export_slides_libreoffice(pptx_path: str, work_dir: str) -> Optional[List[str]]:
    """Utilise LibreOffice pour convertir les slides en PNG.
    Retourne la liste des fichiers générés triés ou None si échec."""
    try:
        cmd = [
            "soffice", "--headless", "--convert-to", "png", pptx_path, "--outdir", work_dir
        ]
        subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
        # LibreOffice génère un .png par slide, nom dérivé
        pngs = sorted(glob.glob(os.path.join(work_dir, "*.png")))
        if pngs:
            return pngs
    except Exception as e:
        print(f"⚠️  Conversion LibreOffice échouée: {e}")
    return None

# ----------------------------------------------------------
# Fallback: créer une image synthétique depuis le texte de la slide
# ----------------------------------------------------------

def build_text_image(text: str, width: int = 1400, padding: int = 32) -> Image.Image:
    lines = []
    words = text.split()
    line = ""
    max_chars = 80
    for w in words:
        candidate = (line + " " + w).strip()
        if len(candidate) > max_chars:
            lines.append(line)
            line = w
        else:
            line = candidate
    if line:
        lines.append(line)
    font = ImageFont.load_default()
    line_height = font.getbbox("Ag")[3] + 6
    height = padding * 2 + line_height * len(lines)
    img = Image.new("RGB", (width, height), "white")
    draw = ImageDraw.Draw(img)
    y = padding
    for l in lines:
        draw.text((padding, y), l, fill="black", font=font)
        y += line_height
    return img

# ----------------------------------------------------------
# Description d'une slide via Pixtral
# ----------------------------------------------------------

def describe_slide(processor, model, image: Image.Image, raw_text: str, max_new_tokens: int = 192) -> str:
    prompt = (
        "Tu es un assistant qui décrit des diapositives de présentation en français de façon structurée. "
        "Analyse l'image fournie et le texte brut (si présent). Fournis une description détaillée incluant: \n"
        "- Le titre s'il est identifiable\n"
        "- Les points clés / sections\n"
        "- Le type de contenu (liste, schéma, tableau, image illustrative...)\n"
        "- L'intention pédagogique\n"
        "Format: un paragraphe synthétique suivi d'une liste à puces claire. Ne pas halluciner."  # Pas d'invention
    )
    # On fournit le texte brut en suffixe pour aider le modèle
    conditioning_text = f"Texte OCR / brut éventuel: {raw_text[:1500]}" if raw_text else "(Pas de texte brut)"

    inputs = processor(images=image, text=f"{prompt}\n\n{conditioning_text}\n\nDescription:", return_tensors="pt").to(model.device)
    with torch.no_grad():
        out_ids = model.generate(**inputs, max_new_tokens=max_new_tokens, do_sample=False)
    decoded = processor.batch_decode(out_ids, skip_special_tokens=True)[0]
    # Essayons de ne garder que la partie après 'Description:' si répétée
    if 'Description:' in decoded:
        decoded = decoded.split('Description:', 1)[1].strip()
    return decoded.strip()

# ----------------------------------------------------------
# Extraction PPTX (texte + images brutes pour fallback)
# ----------------------------------------------------------

def extract_slide_texts(pptx_path: str):
    prs = Presentation(pptx_path)
    slides_meta = []
    for idx, slide in enumerate(prs.slides, start=1):
        title = None
        # Titre: shape avec placeholder titre ou première zone texte non vide
        for shape in slide.shapes:
            if hasattr(shape, 'placeholder_format') and shape.placeholder_format.type == 1:  # TITLE
                if hasattr(shape, 'text') and shape.text.strip():
                    title = shape.text.strip()
                    break
        if title is None:
            for shape in slide.shapes:
                if hasattr(shape, 'has_text_frame') and shape.has_text_frame and shape.text_frame.text.strip():
                    title = shape.text_frame.text.strip().split('\n')[0]
                    break
        # Agréger tout le texte
        all_text_parts = []
        for shape in slide.shapes:
            if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                txt = '\n'.join(p.text for p in shape.text_frame.paragraphs if p.text.strip())
                if txt.strip():
                    all_text_parts.append(txt.strip())
        raw_text = '\n'.join(all_text_parts)
        slides_meta.append({
            'number': idx,
            'title': title or f"Slide {idx}",
            'raw_text': raw_text
        })
    return slides_meta

# ----------------------------------------------------------
# Main
# ----------------------------------------------------------

def parse_args():
    ap = argparse.ArgumentParser(description="Description détaillée des slides via Pixtral")
    ap.add_argument('--pptx', required=True, help='Fichier PPTX')
    ap.add_argument('--out', default='descriptions.xlsx', help='Fichier Excel de sortie')
    ap.add_argument('--tmp', default=None, help='Dossier temporaire (sinon auto)')
    ap.add_argument('--model', default='mistral-community/pixtral-12b', help='Modèle vision Pixtral')
    ap.add_argument('--quant', default='4bit', choices=['4bit','8bit','fp16','auto'], help='Quantization (défaut 4bit)')
    ap.add_argument('--attn', default='eager', choices=['eager','flash_attention_2','sdpa'], help='Implémentation attention (eager plus sûr)')
    ap.add_argument('--no-trust-remote-code', action='store_true', help="Désactive trust_remote_code")
    ap.add_argument('--cpu', action='store_true', help='Forcer exécution CPU (diagnostic)')
    ap.add_argument('--offload-dir', default=None, help='Dossier pour offload (4bit/8bit)')
    ap.add_argument('--max-gpu-mem', default=None, help='Ex: 13GiB pour limiter usage GPU')
    ap.add_argument('--fast-processor', action='store_true', help='Forcer processeur image rapide si dispo')
    ap.add_argument('--max-new', type=int, default=192, help='Max tokens génération')
    ap.add_argument('--no-libreoffice', action='store_true', help='Ne pas tenter d\'utiliser LibreOffice pour rendu PNG')
    return ap.parse_args()


def main():
    args = parse_args()
    if not os.path.exists(args.pptx):
        print(f"❌ Fichier introuvable: {args.pptx}")
        sys.exit(1)

    processor, model = load_pixtral(
        args.model,
        args.quant,
        attn_impl=args.attn,
        trust_remote_code=not args.no_trust_remote_code,
        force_cpu=args.cpu,
        offload_dir=args.offload_dir,
        max_gpu_mem=args.max_gpu_mem,
        fast_processor=args.fast_processor,
    )

    slides_meta = extract_slide_texts(args.pptx)

    work_dir = args.tmp or tempfile.mkdtemp(prefix='pixtral_slides_')
    os.makedirs(work_dir, exist_ok=True)
    print(f"📂 Dossier temp: {work_dir}")

    images = None
    if not args.no_libreoffice:
        images = export_slides_libreoffice(args.pptx, work_dir)
        if images:
            print(f"✅ {len(images)} images rendues via LibreOffice")
        else:
            print("⚠️ Rendu visuel indisponible, fallback images texte")
    else:
        print("ℹ️ LibreOffice désactivé, fallback images texte")

    records = []
    for meta in slides_meta:
        idx = meta['number']
        title = meta['title']
        raw_text = meta['raw_text']
        # Choisir image: si rendu dispo, prendre celle correspondant à l'index (approx)
        slide_img = None
        if images and len(images) >= idx:
            try:
                slide_img = Image.open(images[idx-1]).convert('RGB')
            except Exception:
                slide_img = None
        if slide_img is None:
            slide_img = build_text_image(f"{title}\n\n{raw_text}")

        print(f"🖼️  Slide {idx}: description en cours...")
        description = describe_slide(processor, model, slide_img, raw_text, max_new_tokens=args.max_new)
        records.append({
            'Slide Number': idx,
            'Slide Title': title,
            'description détaillée de la slides': description
        })

    df = pd.DataFrame(records).sort_values('Slide Number')
    df.to_excel(args.out, index=False)
    print(f"✅ Fichier généré: {args.out}")

if __name__ == '__main__':
    main()
