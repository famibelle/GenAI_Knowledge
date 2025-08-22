#!/usr/bin/env python3
"""
Description détaillée de chaque slide d'un PPTX avec un modèle Mistral 7B (texte seulement).
Génère un Excel avec colonnes:
  - Slide Number
  - Slide Title
  - description détaillée de la slides

Le modèle n'étant pas vision, la description repose sur l'extraction structurée du texte (titres, listes, tableaux détectés naïvement).

Prérequis:
  pip install python-pptx pandas transformers accelerate bitsandbytes

Exemple:
  python pptx_mistral7b_describer.py --pptx Pptx/AKABI\ -\ AI\ KT\ -\ Module\ 1\ .pptx --out descriptions_mistral7b.xlsx \
      --model mistral-community/Mistral-7B-Instruct-v0.3 --quant 4bit
"""
from __future__ import annotations
import argparse
import os
import sys
from dataclasses import dataclass
from typing import List, Optional, Dict, Any
import pandas as pd
from pptx import Presentation
try:
    from pptx.enum.shapes import PP_PLACEHOLDER
except Exception:  # compat
    PP_PLACEHOLDER = None
import torch
from transformers import AutoModelForCausalLM, AutoTokenizer
try:
    from transformers import BitsAndBytesConfig
except ImportError:
    BitsAndBytesConfig = None

# ----------------------------- Extraction PPTX -----------------------------
@dataclass
class SlideContent:
    number: int
    title: str
    bullet_blocks: List[List[str]]  # liste de blocs (liste de puces)
    paragraphs: List[str]          # paragraphes libres
    raw_text: str                  # concat global


def extract_pptx(path: str) -> List[SlideContent]:
    prs = Presentation(path)
    slides: List[SlideContent] = []
    for idx, slide in enumerate(prs.slides, start=1):
        title = None
        for shape in slide.shapes:
            # Accès sécurisé au placeholder
            pf_type = None
            try:
                pf = shape.placeholder_format  # peut lever ValueError
                pf_type = getattr(pf, 'type', None)
            except Exception:
                pf = None
            # Types: 1 = TITLE (ou utiliser PP_PLACEHOLDER.TITLE si dispo)
            is_title_placeholder = (pf_type == 1) or (PP_PLACEHOLDER and pf_type == getattr(PP_PLACEHOLDER, 'TITLE', -999))
            if is_title_placeholder:
                txt = getattr(shape, 'text', '')
                if txt and txt.strip():
                    title = txt.strip()
                    break
        bullet_blocks: List[List[str]] = []
        paragraphs: List[str] = []
        for shape in slide.shapes:
            if not getattr(shape, 'has_text_frame', False):
                continue
            tf = shape.text_frame
            current_block: List[str] = []
            for p in tf.paragraphs:
                txt = p.text.strip()
                if not txt:
                    continue
                # Détection simple d'une puce (niveau > 0 ou bullet == True)
                if p.level > 0 or getattr(p, 'bullet', None):
                    current_block.append(txt)
                else:
                    # Si on accumulait une liste, on la pousse
                    if current_block:
                        bullet_blocks.append(current_block)
                        current_block = []
                    paragraphs.append(txt)
            if current_block:
                bullet_blocks.append(current_block)
        # fallback titre
        if not title and paragraphs:
            title = paragraphs[0][:120]
        title = title or f"Slide {idx}"
        # Construire raw
        raw_parts = [f"TITRE: {title}"]
        for block in bullet_blocks:
            raw_parts.append("LISTE: " + " | ".join(block))
        for par in paragraphs:
            raw_parts.append("PARAGRAPHE: " + par)
        raw_text = "\n".join(raw_parts)
        slides.append(SlideContent(idx, title, bullet_blocks, paragraphs, raw_text))
    return slides

# ----------------------------- Modèle Mistral 7B ----------------------------

def load_model(model_name: str, quant: str, max_gpu_mem: Optional[str], force_cpu: bool = False, use_fast: bool = True):
    gpu = torch.cuda.is_available() and not force_cpu
    load_kwargs: Dict[str, Any] = {"device_map": "auto" if not force_cpu else "cpu"}
    quant = (quant or 'auto').lower()

    def bnb_cfg(mode: str):
        if BitsAndBytesConfig is None:
            return None
        if mode == '4bit':
            return BitsAndBytesConfig(load_in_4bit=True, bnb_4bit_quant_type='nf4', bnb_4bit_use_double_quant=True,
                                      bnb_4bit_compute_dtype=torch.bfloat16 if gpu and torch.cuda.is_bf16_supported() else torch.float16)
        if mode == '8bit':
            return BitsAndBytesConfig(load_in_8bit=True)
        return None

    attempt_order = []
    if quant in ('4bit','4'): attempt_order = ['4bit']
    elif quant in ('8bit','8'): attempt_order = ['8bit']
    elif quant in ('fp16','half','16'): attempt_order = ['fp16']
    else: attempt_order = ['4bit','8bit','fp16']

    if max_gpu_mem and gpu:
        load_kwargs['max_memory'] = {0: max_gpu_mem, 'cpu': '48GiB'}

    try:
        tokenizer = AutoTokenizer.from_pretrained(model_name, use_fast=use_fast)
    except ValueError as ve:
        msg = str(ve)
        if 'sentencepiece' in msg.lower():
            raise RuntimeError(
                "Le tokenizer nécessite 'sentencepiece'. Installe: pip install sentencepiece && ré-exécute. "
                "Ex: pip install sentencepiece tokenizers" ) from ve
        # réessayer sans fast
        if use_fast:
            tokenizer = AutoTokenizer.from_pretrained(model_name, use_fast=False)
        else:
            raise
    last_err = None
    for mode in attempt_order:
        local = dict(load_kwargs)
        if mode in ('4bit','8bit') and gpu:
            cfg = bnb_cfg(mode)
            if cfg:
                local['quantization_config'] = cfg
        elif mode == 'fp16' and gpu:
            local['torch_dtype'] = torch.float16
        print(f"⚙️  Tentative chargement {model_name} mode={mode} ...")
        try:
            model = AutoModelForCausalLM.from_pretrained(model_name, **local)
            print(f"✅ Modèle chargé ({mode})")
            return tokenizer, model
        except Exception as e:
            print(f"⚠️  Échec mode {mode}: {e}")
            last_err = e
    print("↪️  Fallback CPU float32 ...")
    model = AutoModelForCausalLM.from_pretrained(model_name, device_map='cpu')
    return tokenizer, model

# ----------------------------- Génération description ----------------------

def build_prompt(slide: SlideContent) -> str:
    # On structure le prompt pour faire produire une description pédagogique
    bullet_section = "".join(
        f"\n- Bloc de puces:\n" + "\n".join(f"  * {item}" for item in block)
        for block in slide.bullet_blocks
    )
    paragraph_section = "".join(f"\nParagraphe: {p}" for p in slide.paragraphs)
    base = f"""Tu es un assistant pédagogique francophone.
Décris de manière détaillée et structurée la diapositive suivante.
Objectif: produire une description claire pour quelqu'un qui ne voit pas la slide.
Inclure:
- Intention pédagogique implicite
- Résumé synthétique
- Détails structurés (listes / sections)
- Concepts clés
- Suggestions d'amélioration si utiles (facultatif, 1 phrase)
Ne pas inventer de contenu absent.

SLIDE #{slide.number}
Titre: {slide.title}
Contenu brut:{bullet_section}{paragraph_section}

Réponse (format: paragraphe initial + liste à puces):
"""
    return base


def generate_description(tokenizer, model, prompt: str, max_new_tokens: int, temperature: float, top_p: float, use_chat: bool = True) -> str:
    input_ids = None
    if use_chat and hasattr(tokenizer, 'apply_chat_template'):
        # Fusionner system + user dans un seul message user pour éviter l'erreur d'alternance stricte.
        merged_user = ("[Contexte système] Tu es un assistant pédagogique précis et factuel.\n\n" + prompt)
        messages = [{"role": "user", "content": merged_user}]
        try:
            model_inputs = tokenizer.apply_chat_template(messages, tokenize=True, return_tensors='pt')
            input_ids = model_inputs.to(model.device)
        except Exception as e:
            print(f"⚠️ Chat template échoue ({e}), fallback prompt brut")
    if input_ids is None:
        input_ids = tokenizer(prompt, return_tensors='pt').input_ids.to(model.device)

    with torch.no_grad():
        gen_ids = model.generate(
            input_ids,
            max_new_tokens=max_new_tokens,
            do_sample=temperature > 0.01,
            temperature=temperature,
            top_p=top_p,
            repetition_penalty=1.05,
            pad_token_id=tokenizer.eos_token_id,
        )
    text = tokenizer.decode(gen_ids[0], skip_special_tokens=True)
    if 'Réponse' in text:
        after = text.split('Réponse',1)[1].strip()
        if after:
            text = after

    print(text)
    return text.strip()

# ----------------------------- Main ----------------------------------------

def parse_args():
    ap = argparse.ArgumentParser(description='Description PPTX via Mistral 7B (texte)')
    ap.add_argument('--pptx', required=True, help='Fichier PPTX')
    ap.add_argument('--out', default='descriptions_mistral7b.xlsx', help='Excel de sortie')
    ap.add_argument('--model', default='mistral-community/Mistral-7B-Instruct-v0.3', help='Modèle Mistral Instruct')
    ap.add_argument('--quant', default='4bit', choices=['4bit','8bit','fp16','auto'], help='Quantization')
    ap.add_argument('--max-gpu-mem', default=None, help='Ex: 12GiB')
    ap.add_argument('--cpu', action='store_true', help='Forcer CPU')
    ap.add_argument('--no-fast-tokenizer', action='store_true', help='Désactive tokenizer rapide')
    ap.add_argument('--max-new', type=int, default=320, help='Max nouveaux tokens')
    ap.add_argument('--temperature', type=float, default=0.4, help='Température')
    ap.add_argument('--top-p', type=float, default=0.9, help='Top-p')
    ap.add_argument('--no-chat', action='store_true', help='Désactive format chat template')
    return ap.parse_args()


def main():
    args = parse_args()
    if not os.path.exists(args.pptx):
        print(f"❌ Fichier introuvable: {args.pptx}")
        sys.exit(1)

    print("📥 Extraction du PPTX ...")
    slides = extract_pptx(args.pptx)
    print(f"✅ {len(slides)} slides extraites")

    print("⚙️ Chargement du modèle ...")
    tokenizer, model = load_model(
        args.model,
        args.quant,
        args.max_gpu_mem,
        force_cpu=args.cpu,
        use_fast=not args.no_fast_tokenizer,
    )

    records = []
    for slide in slides:
        print(f"🧠 Génération description slide {slide.number} ...")
        prompt = build_prompt(slide)
        desc = generate_description(tokenizer, model, prompt, args.max_new, args.temperature, args.top_p, use_chat=not args.no_chat)
        records.append({
            'Slide Number': slide.number,
            'Slide Title': slide.title,
            'description détaillée de la slides': desc
        })

    df = pd.DataFrame(records).sort_values('Slide Number')
    df.to_excel(args.out, index=False)
    print(f"✅ Fichier généré: {args.out}")

if __name__ == '__main__':
    main()
