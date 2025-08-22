#!/usr/bin/env python3
"""
Extraction de texte de slides PowerPoint avec OCR local + post-traitement (optionnel) par un LLM Mistral.

Objectifs:
 1. Parcourir un fichier .pptx
 2. Extraire le texte natif (zones de texte) sans OCR (plus fiable)
 3. Extraire les images et appliquer un OCR local (PaddleOCR ou Tesseract en fallback)
 4. (Optionnel) Nettoyer / restructurer / normaliser le texte via un modèle Mistral local (HF transformers)
 5. Sauvegarder :
    - Un CSV / XLSX récapitulatif (slide, ordre, type, texte)
    - Un dossier par slide avec images extraites

Dépendances possibles (installez selon vos besoins) :
  pip install python-pptx paddleocr opencv-python pandas pillow transformers accelerate bitsandbytes sentencepiece
  (Tesseract en fallback OCR: installer binaire système + pip install pytesseract)

Remarques:
 - Il n'existe pas (à ce jour) de "Mistral OCR" officiel. Ici, on combine un moteur OCR visuel (PaddleOCR/Tesseract) + un modèle Mistral pour le post-traitement linguistique.
 - Le post-traitement Mistral est optionnel (--mistral-clean). Il peut corriger l'espacement, enlever du bruit, reformater.
"""
from __future__ import annotations
import argparse
import os
import io
import sys
import tempfile
import json
from dataclasses import dataclass, asdict
from typing import List, Optional, Dict, Any
import torch  # Ajout pour utilisation dans MistralCleaner et Vision OCR

# --- Imports de base ---
try:
    from pptx import Presentation
except ImportError:
    print("❌ python-pptx requis: pip install python-pptx")
    sys.exit(1)

from PIL import Image
import pandas as pd

# OCR Engines (on essaiera Paddle d'abord, sinon Tesseract si présent)
_PADDLE_OK = False
try:
    from paddleocr import PaddleOCR  # type: ignore
    _PADDLE_OK = True
except Exception:
    pass

_TESSERACT_OK = False
try:
    import pytesseract  # type: ignore
    _TESSERACT_OK = True
except Exception:
    pass

# Mistral (LLM) pour post-traitement optionnel
_MISTRAL_OK = False
try:
    from transformers import AutoModelForCausalLM, AutoTokenizer, TextGenerationPipeline
    _MISTRAL_OK = True
except Exception:
    pass

# --- Vision OCR (Mistral multimodal / Pixtral) ---
try:
    from transformers import AutoProcessor, AutoModelForCausalLM  # multimodal
    _MISTRAL_VISION_OK = True
except Exception:
    _MISTRAL_VISION_OK = False

class MistralVisionOCREngine:
    """OCR via modèle multimodal (ex: Pixtral) avec support quantization (4-bit / 8-bit / fp16)."""
    def __init__(self, model_name: str, device: Optional[str] = None, max_new_tokens: int = 256, quant_mode: str = "4bit"):
        if not _MISTRAL_VISION_OK:
            raise RuntimeError("transformers multimodal indisponible pour vision OCR.")
        self.model_name = model_name
        self.max_new_tokens = max_new_tokens
        self.device = device or ("cuda" if torch.cuda.is_available() else "cpu")
        self.quant_mode = quant_mode.lower()
        print(f"🤖 Chargement modèle vision: {model_name} (quant={self.quant_mode})")

        load_kwargs = {"device_map": "auto"}

        # Sélection dtype compute (T4 ne supporte pas bfloat16)
        gpu_fp16 = torch.cuda.is_available()
        compute_dtype = torch.float16 if gpu_fp16 else torch.float32

        try:
            if self.quant_mode in ("4bit", "4", "nf4"):
                load_kwargs.update({
                    "load_in_4bit": True,
                    "bnb_4bit_quant_type": "nf4",
                    "bnb_4bit_use_double_quant": True,
                    "bnb_4bit_compute_dtype": compute_dtype
                })
            elif self.quant_mode in ("8bit", "8"):
                load_kwargs.update({"load_in_8bit": True})
            elif self.quant_mode in ("fp16", "half", "16") and gpu_fp16:
                load_kwargs.update({"torch_dtype": torch.float16})
            else:
                # défaut: laisser AutoModel choisir (fp32 / bf16 si support)
                if gpu_fp16:
                    load_kwargs.setdefault("torch_dtype", compute_dtype)
        except Exception as e:
            print(f"⚠️ Quantization non appliquée ({e}), fallback plein format.")

        self.processor = AutoProcessor.from_pretrained(model_name)
        self.model = AutoModelForCausalLM.from_pretrained(model_name, **load_kwargs)
        print("✅ Modèle vision chargé")

    def ocr(self, image_path: str) -> str:
        from PIL import Image as _Img
        try:
            img = _Img.open(image_path).convert("RGB")
            prompt = ("Transcris fidèlement le texte lisible (français) de cette image de slide. "
                      "Ne pas inventer. Ne pas traduire. Sortie: texte brut.")
            inputs = self.processor(images=img, text=prompt, return_tensors="pt").to(self.model.device)
            out_ids = self.model.generate(**inputs, max_new_tokens=self.max_new_tokens, do_sample=False)
            decoded = self.processor.batch_decode(out_ids, skip_special_tokens=True)[0]
            # Retirer le prompt si réinséré
            if prompt in decoded:
                decoded = decoded.split(prompt, 1)[-1].strip()
            return basic_cleanup(decoded)
        except Exception as e:
            return f"[VISION_OCR_ERROR:{e}]"

@dataclass
class ExtractedItem:
    slide_index: int
    slide_id: str
    order: int
    kind: str            # 'text' ou 'ocr-image'
    raw_text: str
    cleaned_text: Optional[str] = None
    source_image: Optional[str] = None

# -------------------------------------------------------------
# Post-traitement (nettoyage basique avant Mistral)
# -------------------------------------------------------------
def basic_cleanup(text: str) -> str:
    import re
    t = text.replace('\r', ' ').replace('\n', ' ')
    t = re.sub(r'\s+', ' ', t).strip()
    return t

# -------------------------------------------------------------
# OCR wrappers
# -------------------------------------------------------------
class OCREngine:
    def __init__(self, lang: str = 'fr'):  # Adapter la langue
        self.lang = lang
        self._engine = None
        self.engine_name = None
        if _PADDLE_OK:
            # PaddleOCR: multi-lang possible, 'fr' nécessite le modèle language-specific ou use_angle_cls
            try:
                self._engine = PaddleOCR(lang='fr', show_log=False)
                self.engine_name = 'paddleocr'
            except Exception:
                self._engine = None
        if self._engine is None and _TESSERACT_OK:
            self.engine_name = 'tesseract'

    def available(self) -> bool:
        return self.engine_name is not None

    def ocr_image(self, img_path: str) -> str:
        if self.engine_name == 'paddleocr' and self._engine is not None:
            result = self._engine.ocr(img_path, cls=False)
            lines = []
            for block in result:
                if not block:
                    continue
                for line in block:
                    if len(line) >= 2 and isinstance(line[1], (list, tuple)):
                        text_part = line[1][0]
                    else:
                        text_part = line[1][0] if isinstance(line[1], (list, tuple)) else ''
                    if text_part:
                        lines.append(text_part)
            return basic_cleanup('\n'.join(lines))
        elif self.engine_name == 'tesseract':
            return basic_cleanup(pytesseract.image_to_string(Image.open(img_path), lang=self.lang))
        else:
            return ''

# -------------------------------------------------------------
# Mistral Post-processing
# -------------------------------------------------------------
class MistralCleaner:
    def __init__(self, model_name: str, device: Optional[str] = None, load_8bit: bool = True):
        if not _MISTRAL_OK:
            raise RuntimeError("transformers non installé (pip install transformers accelerate bitsandbytes sentencepiece)")
        self.model_name = model_name
        self.device = device or ("cuda" if torch.cuda.is_available() else "cpu")  # type: ignore
        from transformers import AutoModelForCausalLM, AutoTokenizer
        load_kwargs = {}
        if load_8bit:
            try:
                load_kwargs.update({"load_in_8bit": True, "device_map": "auto"})
            except Exception:
                pass
        self.tokenizer = AutoTokenizer.from_pretrained(model_name)
        self.model = AutoModelForCausalLM.from_pretrained(model_name, **load_kwargs)

    def clean_text(self, raw: str, max_new_tokens: int = 128) -> str:
        prompt = (
            "Nettoie et reformate le texte OCR suivant (français) en phrases correctes, "
            "sans ajouter d'informations nouvelles. Conserve les listes si pertinentes.\n\nTexte: "
            f"""{raw}\n\nTexte nettoyé:"""
        )
        inputs = self.tokenizer(prompt, return_tensors='pt')
        if hasattr(self.model, 'device'):
            inputs = {k: v.to(self.model.device) for k, v in inputs.items()}
        output_ids = self.model.generate(**inputs, max_new_tokens=max_new_tokens, do_sample=False)
        generated = self.tokenizer.decode(output_ids[0], skip_special_tokens=True)
        # Extraire la partie après "Texte nettoyé:" si elle existe
        if "Texte nettoyé:" in generated:
            cleaned = generated.split("Texte nettoyé:", 1)[1].strip()
        else:
            cleaned = generated.strip()
        return cleaned

# -------------------------------------------------------------
# Extraction PPTX
# -------------------------------------------------------------
class PPTXExtractor:
    def __init__(self, path: str, ocr_engine: Optional[OCREngine] = None, out_dir: str = "pptx_ocr_output"):
        self.path = path
        self.ocr_engine = ocr_engine
        self.out_dir = out_dir
        os.makedirs(out_dir, exist_ok=True)

    def extract(self) -> List[ExtractedItem]:
        prs = Presentation(self.path)
        all_items: List[ExtractedItem] = []
        for idx, slide in enumerate(prs.slides, start=1):
            slide_dir = os.path.join(self.out_dir, f"slide_{idx:03d}")
            os.makedirs(slide_dir, exist_ok=True)
            order = 0
            # Text shapes
            for shape in slide.shapes:
                if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                    text = '\n'.join(p.text for p in shape.text_frame.paragraphs)
                    order += 1
                    all_items.append(ExtractedItem(
                        slide_index=idx,
                        slide_id=f"slide_{idx:03d}",
                        order=order,
                        kind='text',
                        raw_text=basic_cleanup(text)
                    ))
                # Images pour OCR
                if shape.shape_type == 13 or getattr(shape, 'image', None) is not None:  # 13 = PICTURE (valeur interne)
                    try:
                        image = shape.image
                        image_bytes = image.blob
                        image_ext = image.ext or 'png'
                        img_filename = os.path.join(slide_dir, f"img_{order+1}.{image_ext}")
                        with open(img_filename, 'wb') as f:
                            f.write(image_bytes)
                        order += 1
                        ocr_text = ''
                        if self.ocr_engine and self.ocr_engine.available():
                            ocr_text = self.ocr_engine.ocr_image(img_filename)
                        all_items.append(ExtractedItem(
                            slide_index=idx,
                            slide_id=f"slide_{idx:03d}",
                            order=order,
                            kind='ocr-image',
                            raw_text=ocr_text,
                            source_image=img_filename
                        ))
                    except Exception:
                        pass
        return all_items

# -------------------------------------------------------------
# Sauvegarde
# -------------------------------------------------------------
class ResultWriter:
    def __init__(self, items: List[ExtractedItem], out_dir: str):
        self.items = items
        self.out_dir = out_dir

    def to_csv(self):
        df = pd.DataFrame([asdict(it) for it in self.items])
        path = os.path.join(self.out_dir, 'extraction.csv')
        df.to_csv(path, index=False)
        return path

    def to_excel(self):
        df = pd.DataFrame([asdict(it) for it in self.items])
        path = os.path.join(self.out_dir, 'extraction.xlsx')
        df.to_excel(path, index=False)
        return path

    def to_json(self):
        path = os.path.join(self.out_dir, 'extraction.json')
        with open(path, 'w', encoding='utf-8') as f:
            json.dump([asdict(it) for it in self.items], f, ensure_ascii=False, indent=2)
        return path

# -------------------------------------------------------------
# Main
# -------------------------------------------------------------
def parse_args():
    p = argparse.ArgumentParser(description="Extraction PPTX + OCR + (optionnel) nettoyage Mistral")
    p.add_argument('--pptx', required=True, help='Fichier .pptx en entrée')
    p.add_argument('--out', default='pptx_ocr_output', help='Dossier de sortie')
    p.add_argument('--lang', default='fr', help='Langue OCR (fr, en, etc.)')
    p.add_argument('--no-ocr', action='store_true', help='Désactiver totalement l\'OCR (texte natif uniquement)')
    p.add_argument('--mistral-clean', action='store_true', help='Activer le post-traitement Mistral (texte)')
    p.add_argument('--mistral-model', default='mistralai/Mistral-7B-Instruct-v0.2', help='Nom du modèle HF Mistral (texte)')
    p.add_argument('--no-8bit', action='store_true', help='Ne pas charger le modèle Mistral texte en 8-bit')
    p.add_argument('--excel', action='store_true', help='Exporter aussi en Excel')
    p.add_argument('--vision-ocr', action='store_true', help='Utiliser un modèle vision Mistral (ex pixtral) pour OCR des images')
    p.add_argument('--vision-model', default='mistral-community/pixtral-12b', help='Modèle vision (multimodal) Mistral à utiliser')
    p.add_argument('--vision-max-new', type=int, default=192, help='Max tokens générés pour vision OCR')
    p.add_argument('--vision-quant', default='4bit', choices=['4bit','8bit','fp16','auto'], help='Mode de quantization du modèle vision (défaut: 4bit)')
    return p.parse_args()


def main():
    args = parse_args()

    if not os.path.exists(args.pptx):
        print(f"❌ Fichier introuvable: {args.pptx}")
        sys.exit(1)

    print(f"📁 Fichier PPTX: {args.pptx}")
    print(f"📂 Dossier sortie: {args.out}")

    # OCR engine traditionnel (Paddle/Tesseract) si demandé
    ocr_engine = None
    if not args.no_ocr and not args.vision_ocr:
        ocr_engine = OCREngine(lang=args.lang)
        if not (ocr_engine and ocr_engine.available()):
            print("⚠️ Aucun moteur OCR disponible (paddle/tesseract). Mode texte-natif uniquement.")
            ocr_engine = None
        else:
            print(f"🔎 OCR classique actif via: {ocr_engine.engine_name}")
    elif args.vision_ocr:
        print("🖼️ OCR vision Mistral sélectionné (les images seront traitées après extraction)…")
    else:
        print("ℹ️ OCR désactivé (no-ocr)")

    extractor = PPTXExtractor(args.pptx, ocr_engine=ocr_engine, out_dir=args.out)
    items = extractor.extract()
    print(f"✅ Extraction brute: {len(items)} éléments")

    # Post-traitement Mistral
    if args.mistral_clean:
        if not _MISTRAL_OK:
            print("❌ transformers non disponible. Installez-le ou retirez --mistral-clean")
        else:
            print("🤖 Chargement du modèle Mistral pour nettoyage...")
            try:
                cleaner = MistralCleaner(args.mistral_model, load_8bit=not args.no_8bit)
                for it in items:
                    if it.raw_text and it.raw_text.strip():
                        try:
                            it.cleaned_text = cleaner.clean_text(it.raw_text)
                        except Exception as e:
                            print(f"⚠️ Nettoyage Mistral échoué (slide {it.slide_index}, ordre {it.order}): {e}")
            except Exception as e:
                print(f"❌ Impossible de charger Mistral: {e}")

    # Si vision OCR demandé, post-traiter les items d'images
    if args.vision_ocr:
        if not _MISTRAL_VISION_OK:
            print("❌ Vision OCR indisponible: installez transformers >= 4.40 + modèle multimodal.")
        else:
            try:
                vision_engine = MistralVisionOCREngine(
                    args.vision_model,
                    max_new_tokens=args.vision_max_new,
                    quant_mode=args.vision_quant
                )
                updated = 0
                for it in items:
                    if it.kind == 'ocr-image' and (not it.raw_text or it.raw_text.strip() == '') and it.source_image:
                        txt = vision_engine.ocr(it.source_image)
                        it.raw_text = txt
                        updated += 1
                print(f"🖼️ Vision OCR effectué sur {updated} images vides.")
            except Exception as e:
                print(f"❌ Échec Vision OCR: {e}")

    writer = ResultWriter(items, args.out)
    csv_path = writer.to_csv()
    json_path = writer.to_json()
    print(f"📄 Export CSV: {csv_path}")
    print(f"📄 Export JSON: {json_path}")
    if args.excel:
        xlsx_path = writer.to_excel()
        print(f"📄 Export Excel: {xlsx_path}")

    print("🎉 Terminé.")

if __name__ == '__main__':
    main()
